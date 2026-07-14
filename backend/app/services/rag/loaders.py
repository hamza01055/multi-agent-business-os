"""Turn any supported source into plain text + metadata."""
from pathlib import Path


def load_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def load_docx(path: str) -> str:
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def load_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


LOADERS = {".pdf": load_pdf, ".docx": load_docx, ".pptx": load_pptx}


def load_file(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext not in LOADERS:
        raise ValueError(f"Unsupported file type: {ext}")
    return LOADERS[ext](path)
